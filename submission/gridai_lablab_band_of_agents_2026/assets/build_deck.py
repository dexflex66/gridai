"""Build GridAI pitch deck for lablab.ai Band of Agents Hackathon submission."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BG = RGBColor(0x0F, 0x19, 0x23)
AMBER = RGBColor(0xF4, 0xA6, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
STEEL = RGBColor(0xB8, 0xC8, 0xD8)
DARK_CARD = RGBColor(0x16, 0x25, 0x35)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x27, 0xAE, 0x60)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # blank layout

ASSETS = "/Users/mayank/gridai/viz/audit_screenshots"
OUT = "/Users/mayank/gridai/submission/gridai_lablab_band_of_agents_2026/assets/gridai_pitch_deck.pptx"


def new_slide():
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def title_bar(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8.8), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = STEEL
    p.font.name = FONT


def add_body(slide, x, y, w, h, runs, font_size=Pt(14)):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for run_spec in runs:
        is_break = run_spec.get("break", False)
        if is_break and p.text != "":
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = run_spec["text"]
        run.font.size = run_spec.get("size", font_size)
        run.font.bold = run_spec.get("bold", False)
        run.font.italic = run_spec.get("italic", False)
        run.font.color.rgb = run_spec.get("color", WHITE)
        run.font.name = FONT
        if "bullet" in run_spec:
            p.level = 0
    return tf


def add_text_box(slide, x, y, w, h, text, size=Pt(14), color=WHITE, bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    return tb


def add_card(slide, x, y, w, h, fill_color=DARK_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1 — Title
# ============================================================
s = new_slide()

add_text_box(s, 0.8, 1.6, 8.0, 0.9, "GridAI", Pt(48), WHITE, True)

add_text_box(s, 0.8, 2.35, 4.5, 0.5, "Anti-Herding Compliance for Distributed Energy Resources", Pt(18), STEEL)
add_text_box(s, 0.8, 3.0, 4.5, 0.4, "Four agents on Band  \u00b7  Regulated & High-Stakes Workflows", Pt(13), AMBER)
add_text_box(s, 0.8, 5.8, 4.5, 0.35, "lablab.ai Band of Agents Hackathon  \u00b7  June 2026", Pt(10), STEEL, italic=True)

# ============================================================
# SLIDE 2 — The Problem
# ============================================================
s = new_slide()
title_bar(s, "THE PROBLEM")

add_text_box(s, 0.6, 1.2, 8.8, 0.5,
    "Australia has 15 GWh of home batteries reading the same price signal.",
    Pt(20), WHITE, True)

lines = [
    "When they all respond at once, they synchronise.",
    "The fleet creates a new evening demand spike instead of smoothing the old one.",
    "471 voltage breaches at the network edge.",
    "This is the herding problem \u2014 and it gets worse as VPPs scale.",
]
y = 1.9
for i, line_text in enumerate(lines):
    color = AMBER if "471" in line_text else WHITE
    bold = "471" in line_text
    add_text_box(s, 0.6, y + i * 0.35, 5.5, 0.35, line_text, Pt(14), color, bold)

# Visual: batteries → spike diagram
# 3 battery shapes
for i in range(3):
    bx = 6.5 + i * 0.5
    bat = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(2.8), Inches(0.35), Inches(0.5))
    bat.fill.solid()
    bat.fill.fore_color.rgb = AMBER
    bat.line.fill.background()
add_text_box(s, 6.5, 3.35, 1.8, 0.25, "batteries", Pt(9), STEEL, align=PP_ALIGN.CENTER)

# arrow
add_text_box(s, 8.3, 2.7, 0.6, 0.4, "\u2192", Pt(24), STEEL, align=PP_ALIGN.CENTER)

# "same signal" label
add_card(s, 6.8, 2.3, 2.0, 0.3)
add_text_box(s, 6.8, 2.3, 2.0, 0.3, "same price signal", Pt(9), WHITE, align=PP_ALIGN.CENTER)

# synchronised spike rectangle
spike = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), Inches(2.15), Inches(0.7), Inches(1.5))
spike.fill.solid()
spike.fill.fore_color.rgb = RED
spike.line.fill.background()
add_text_box(s, 9.0, 3.7, 0.7, 0.5, "sync\nspike", Pt(8), RED, align=PP_ALIGN.CENTER)

add_text_box(s, 0.6, 4.6, 8.8, 0.4,
    "Current DERMS/VPP platforms (EnergyHub, Kraken, AutoGrid) coordinate DERs \u2014 none expose anti-herding assurance.",
    Pt(12), STEEL, italic=True)

# ============================================================
# SLIDE 3 — Why It's Hard
# ============================================================
s = new_slide()
title_bar(s, "WHY IT'S HARD")

add_text_box(s, 0.6, 1.2, 8.8, 0.45,
    "The problem is second-order: it appears precisely when VPP coordination succeeds.",
    Pt(16), WHITE, True)
add_text_box(s, 0.6, 1.65, 8.8, 0.35,
    "Each aggregator's fleet looks fine in isolation. The failure is at the system level.",
    Pt(12), STEEL)

# Two-column table
exists = ["DERMS platforms coordinate DERs well", "Voltage constraint management is table stakes", "Real-time dispatch of storage fleets"]
missing = ["Anti-herding assurance before dispatch", "Cause-attributed breach compliance", "Regulator-ready audit trails"]
col_w = 4.3
tx, ty = 0.6, 2.4
rh = 0.55

add_rect(s, tx, ty, col_w, 0.4, DARK_CARD)
add_text_box(s, tx, ty, col_w, 0.4, "WHAT EXISTS", Pt(11), GREEN, True, align=PP_ALIGN.CENTER)

mx = tx + col_w + 0.2
add_rect(s, mx, ty, col_w, 0.4, DARK_CARD)
add_text_box(s, mx, ty, col_w, 0.4, "WHAT'S MISSING", Pt(11), RED, True, align=PP_ALIGN.CENTER)

for i in range(3):
    add_rect(s, tx, ty + 0.4 + i * rh, col_w, rh, DARK_CARD, RGBColor(0x1A, 0x2A, 0x3C))
    add_text_box(s, tx + 0.15, ty + 0.4 + i * rh, col_w - 0.3, rh, exists[i], Pt(10), STEEL, align=PP_ALIGN.LEFT)

    add_rect(s, mx, ty + 0.4 + i * rh, col_w, rh, DARK_CARD, RGBColor(0x1A, 0x2A, 0x3C))
    add_text_box(s, mx + 0.15, ty + 0.4 + i * rh, col_w - 0.3, rh, missing[i], Pt(10), WHITE, align=PP_ALIGN.LEFT)

add_text_box(s, 0.6, 4.5, 8.8, 0.4,
    "Commercial DERMS expose coordination. None expose anti-herding compliance as a primary artifact.",
    Pt(12), STEEL, italic=True)

# ============================================================
# SLIDE 4 — The Mechanism Insight
# ============================================================
s = new_slide()
title_bar(s, "THE MECHANISM INSIGHT")

add_text_box(s, 0.6, 1.2, 8.8, 0.45,
    "Desynchronisation is not just a protocol problem.",
    Pt(16), WHITE, True)
add_text_box(s, 0.6, 1.65, 8.8, 0.35,
    "It depends on fleet-level value heterogeneity. The market design matters as much as the protocol.",
    Pt(12), STEEL)

cards = [
    ("NAIVE", "1.000", "60/60 homes\nsimultaneous", RED),
    ("GOSSIP\nHOMOGENEOUS", "0.367", "22/60 homes\nidentical fleet", AMBER),
    ("GOSSIP\nHETEROGENEOUS", "0.167", "10/60 homes\nrealistic fleet", GREEN),
]
cw, ch = 2.7, 1.8
cy = 2.4
gap = 0.25

for i, (lbl, val, sub, clr) in enumerate(cards):
    cx = 0.6 + i * (cw + gap)
    add_card(s, cx, cy, cw, ch)
    add_text_box(s, cx, cy + 0.12, cw, 0.45, lbl, Pt(10), STEEL, align=PP_ALIGN.CENTER)
    add_text_box(s, cx, cy + 0.55, cw, 0.55, val, Pt(34), clr, True, align=PP_ALIGN.CENTER)
    add_text_box(s, cx, cy + 1.15, cw, 0.45, sub, Pt(9), STEEL, align=PP_ALIGN.CENTER)

add_text_box(s, 0.6, 4.55, 8.8, 0.35,
    "Heterogeneous fleet synchrony is ~55% lower than homogeneous. Market incentives shape the outcome.",
    Pt(12), STEEL, italic=True)

# ============================================================
# SLIDE 5 — The Solution
# ============================================================
s = new_slide()
title_bar(s, "THE SOLUTION")

add_text_box(s, 0.6, 1.2, 5.2, 0.55,
    "GridAI: priority-based coordination protocol with a Band-native compliance layer.",
    Pt(16), WHITE, True)

bullets = [
    ("The Coordinator allocates dispatch slots from global fleet state.", Pt(13), WHITE),
    ("Converges in 1 round.", Pt(13), WHITE),
]
ty = 2.0
for bl, sz, clr in bullets:
    add_text_box(s, 0.8, ty, 5.0, 0.5, "\u2022  " + bl, sz, clr)
    ty += 0.55

ty += 0.2
add_text_box(s, 0.6, ty, 5.0, 0.35, "Battery-herding overvoltage:", Pt(14), WHITE, True)
ty += 0.4

# Big number
add_text_box(s, 0.6, ty, 5.0, 0.55, "471 \u2192 0", Pt(36), AMBER, True)
ty += 0.6

add_text_box(s, 0.6, ty, 5.0, 0.35,
    "Honest tradeoff: residual far-feeder undervoltage, 435 events (distinct from herding, surfaced not hidden)",
    Pt(10), STEEL, italic=True)

img = os.path.join(ASSETS, "03_naive_peak_HERO.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(6.0), Inches(1.2), Inches(3.5), Inches(3.5))

# ============================================================
# SLIDE 6 — The Band Architecture
# ============================================================
s = new_slide()
title_bar(s, "THE BAND ARCHITECTURE")

add_text_box(s, 0.6, 1.15, 8.8, 0.4,
    "Four agents coordinating through Band as the actual collaboration layer.",
    Pt(14), WHITE)

agents = [
    ("FORECASTER", "Identifies risk windows,\nhands off to Coordinator", RGBColor(0x4A, 0x90, 0xD9)),
    ("COORDINATOR", "Runs gossip protocol, hands\ndispatch plan to Compliance", RGBColor(0x50, 0xB8, 0x7D)),
    ("COMPLIANCE", "Checks AS IEC 60038:2022,\nattributes breach cause, escalates", AMBER),
    ("OPERATOR", "Human-in-the-loop, receives\nescalation, records decision", RGBColor(0xE7, 0x4C, 0x3C)),
]

bw, bh = 2.0, 1.5
sx, sy = 0.6, 1.8
gap = 0.3

for i, (name, role, color) in enumerate(agents):
    ax = sx + i * (bw + gap)
    add_card(s, ax, sy, bw, bh)
    add_text_box(s, ax, sy + 0.12, bw, 0.35, name, Pt(11), color, True, align=PP_ALIGN.CENTER)
    add_text_box(s, ax + 0.1, sy + 0.55, bw - 0.2, 0.8, role, Pt(9), STEEL)

    if i < len(agents) - 1:
        arrow_x = ax + bw
        add_text_box(s, arrow_x + 0.02, sy + 0.35, gap - 0.04, 0.4, "\u2192", Pt(20), AMBER, align=PP_ALIGN.CENTER)

# Band layer
add_card(s, 0.6, 3.7, 8.8, 0.35)
add_text_box(s, 0.6, 3.7, 8.8, 0.35,
    "BAND: shared room  \u00b7  @mention routing  \u00b7  append-only audit log  \u00b7  identity & handoff",
    Pt(10), AMBER, align=PP_ALIGN.CENTER)

add_card(s, 1.5, 4.3, 7.0, 0.45)
add_text_box(s, 1.5, 4.3, 7.0, 0.45,
    "Every handoff is traceable. Removing any agent breaks the chain (verified).",
    Pt(12), GREEN, True, align=PP_ALIGN.CENTER)

add_text_box(s, 0.6, 5.15, 8.8, 0.35,
    "Band agents:  @s4142972/gridai-forecaster  \u00b7  coordinator  \u00b7  compliance  \u00b7  operator",
    Pt(9), STEEL, italic=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7 — The Evidence
# ============================================================
s = new_slide()
title_bar(s, "THE EVIDENCE")

rows = [
    [("Metric", STEEL, True), ("Naive", STEEL, True), ("GridAI", STEEL, True)],
    [("Battery-herding overvoltage", WHITE, False), ("471 steps", RED, True), ("0 (eliminated)", GREEN, True)],
    [("Synchrony", WHITE, False), ("1.000", RED, False), ("0.167", GREEN, False)],
    [("Max simultaneous discharge", WHITE, False), ("60/60", RED, False), ("10/60", GREEN, False)],
    [("Convergence", WHITE, False), ("\u2014", STEEL, False), ("1 round", WHITE, False)],
    [("Peak demand reduction", WHITE, False), ("\u2014", STEEL, False), ("\u22120.9% (honest)", STEEL, False)],
]

tx, ty = 0.7, 1.3
col_w = [4.0, 2.3, 2.3]
rh = 0.38

for ri, row in enumerate(rows):
    cy = ty + ri * (rh + 0.02)
    cx = tx
    for ci, (text, color, bold) in enumerate(row):
        cw = col_w[ci]
        if ri == 0:
            add_rect(s, cx, cy, cw, rh, DARK_CARD)
        elif ri % 2 == 0:
            add_rect(s, cx, cy, cw, rh, RGBColor(0x12, 0x20, 0x30))
        add_text_box(s, cx + 0.15, cy, cw - 0.3, rh, text, Pt(11), color, bold, align=PP_ALIGN.LEFT)
        cx += cw + 0.05

add_text_box(s, 0.7, 3.9, 8.6, 0.3,
    "89 tests passing  \u00b7  including provenance coherence & agent interdependence",
    Pt(11), AMBER, True)
add_text_box(s, 0.7, 4.25, 8.6, 0.3,
    "AEMO 2012 Victorian data, 17,568 rows. Representative day: 2012-01-24 (highest evening peak, 8,864 MW).",
    Pt(10), STEEL, italic=True)

# ============================================================
# SLIDE 8 — The Compliance Artifact
# ============================================================
s = new_slide()
title_bar(s, "THE COMPLIANCE ARTIFACT")

img = os.path.join(ASSETS, "06_gossip_compliance_card.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(0.6), Inches(1.2), Inches(5.0), Inches(3.6))

add_text_box(s, 5.9, 1.2, 3.8, 0.55,
    "The Band audit trail is the regulated-workflows deliverable.",
    Pt(15), WHITE, True)

add_text_box(s, 5.9, 1.9, 3.8, 0.3,
    "Every compliance decision is traceable to:",
    Pt(11), STEEL)

trace = [
    "The agent that made it",
    "The data it saw",
    "The moment it happened",
    "The cause category (pv_export vs battery_herding)",
]
for i, item in enumerate(trace):
    add_text_box(s, 6.1, 2.3 + i * 0.32, 3.5, 0.28, "\u2022  " + item, Pt(11), WHITE)

add_text_box(s, 5.9, 3.9, 3.8, 0.5,
    "This is what a network operator or regulator can actually use.",
    Pt(12), STEEL, italic=True)

add_text_box(s, 0.6, 5.3, 8.8, 0.3,
    "AS IEC 60038:2022  \u00b7  CSIP-AUS compliant  \u00b7  8-step Band audit trail per scenario",
    Pt(9), STEEL, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9 — Wedge Positioning
# ============================================================
s = new_slide()
title_bar(s, "WEDGE POSITIONING")

add_text_box(s, 0.6, 1.2, 8.8, 0.45,
    "GridAI is not a replacement for DERMS or VPP platforms.",
    Pt(16), WHITE, True)
add_text_box(s, 0.6, 1.65, 8.8, 0.35,
    "It is an assurance and attribution layer that sits between fleets and the network.",
    Pt(13), STEEL)

layers = [
    ("VPP / DERMS PLATFORMS", "EnergyHub \u00b7 Kraken \u00b7 AutoGrid \u00b7 Tesla VPP", RGBColor(0x4A, 0x90, 0xD9), 2.4),
    ("GRIDAI  \u2190  assurance & attribution", "Anti-herding + cause-attributed compliance", AMBER, 3.05),
    ("DISTRIBUTION NETWORK / REGULATOR", "AEMO \u00b7 SA Power Networks \u00b7 AER", RGBColor(0x27, 0xAE, 0x60), 3.7),
]
lw, lh = 4.0, 0.5
lx = 3.1

for label, sub, color, ly in layers:
    add_card(s, lx, ly, lw, lh)
    add_rect(s, lx, ly, 0.05, lh, color)
    add_text_box(s, lx + 0.2, ly + 0.03, lw - 0.3, 0.25, label, Pt(11), color, True)
    add_text_box(s, lx + 0.2, ly + 0.28, lw - 0.3, 0.2, sub, Pt(8), STEEL)

# arrows
add_text_box(s, lx + lw / 2 - 0.15, 2.9, 0.3, 0.2, "\u2193", Pt(14), STEEL, align=PP_ALIGN.CENTER)
add_text_box(s, lx + lw / 2 - 0.15, 3.55, 0.3, 0.2, "\u2193", Pt(14), STEEL, align=PP_ALIGN.CENTER)

add_text_box(s, 0.6, 4.55, 2.0, 0.3, "NEXT STEPS", Pt(10), STEEL, True)
steps = [
    "\u2022  OpenDSS validation",
    "\u2022  Multi-aggregator scenarios",
    "\u2022  CSIP-AUS live interface",
    "\u2022  RAISE Summit Paris, July 2026",
]
for i, step in enumerate(steps):
    color = AMBER if "RAISE" in step else STEEL
    add_text_box(s, 0.6, 4.85 + i * 0.28, 8.8, 0.28, step, Pt(10), color)

# ============================================================
# SLIDE 10 — Links and Repo
# ============================================================
s = new_slide()

add_text_box(s, 0.6, 0.9, 8.8, 0.55, "Links & Repository", Pt(28), WHITE, True)

links = [
    ("PUBLIC DEMO", "https://dexflex66.github.io/gridai/"),
    ("GITHUB", "https://github.com/dexflex66/gridai"),
    ("BAND AGENTS", "@s4142972/gridai-forecaster · coordinator · compliance · operator"),
]

for i, (label, value) in enumerate(links):
    ly = 1.9 + i * 0.55
    add_card(s, 0.6, ly, 4.5, 0.45)
    add_text_box(s, 0.8, ly, 1.5, 0.45, label, Pt(10), AMBER, True)
    add_text_box(s, 2.3, ly, 2.7, 0.45, value, Pt(11), WHITE)

add_card(s, 0.6, 3.75, 8.8, 1.15)
techs = (
    "89 tests  \u00b7  82 original + 7 causal-link & regression additions\n\n"
    "Built with:  Python  \u00b7  Band SDK  \u00b7  AEMO open data  \u00b7  AS IEC 60038:2022\n"
    "Heterogeneous gossip 0.167  \u2192  Homogeneous gossip 0.367  \u2192  Naive 1.000 synchrony baseline"
)
add_text_box(s, 0.9, 3.9, 8.2, 0.85, techs, Pt(11), STEEL)

add_text_box(s, 0.6, 5.3, 8.8, 0.3,
    "lablab.ai Band of Agents Hackathon  \u00b7  Track: Regulated & High-Stakes Workflows  \u00b7  June 2026",
    Pt(9), STEEL, align=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
prs.save(OUT)
size_kb = os.path.getsize(OUT) / 1024
print(f"Saved: {OUT} ({size_kb:.0f} KB)")